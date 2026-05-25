from sqlalchemy.orm import Session
from app.agents.base import BaseAgent


class PresentationAgent(BaseAgent):
    id = "presentation_agent"
    name = "Presentation Agent"
    agent_type = "presentation"
    description = "Creates presentation outlines, slide structures, and content plans"
    risk_level = "low"
    requires_approval_for = []

    async def run(self, input_data: dict, db: Session) -> dict:
        topic = input_data.get("topic", input_data.get("title", "presentation"))
        context = input_data.get("context", "")

        prompt = f"""Create a presentation outline for: "{topic}"
Additional context: {context}

Return JSON only:
{{
  "title": "presentation title",
  "objective": "what this presentation achieves",
  "audience": "target audience",
  "duration_minutes": 30,
  "slides": [
    {{
      "slide_number": 1,
      "title": "slide title",
      "type": "title|content|data|summary",
      "key_points": ["point 1", "point 2"],
      "notes": "speaker notes or data needed"
    }}
  ],
  "data_needed": ["data point 1", "data point 2"],
  "next_steps": ["action 1", "action 2"]
}}"""
        result = await self.ollama.classify_json(prompt)
        if not result:
            return {
                "title": f"Presentation: {topic}",
                "objective": "To be defined",
                "audience": "To be defined",
                "duration_minutes": 30,
                "slides": [
                    {"slide_number": 1, "title": "Introduction", "type": "title", "key_points": [topic], "notes": ""},
                    {"slide_number": 2, "title": "Overview", "type": "content", "key_points": [], "notes": ""},
                    {"slide_number": 3, "title": "Key Points", "type": "content", "key_points": [], "notes": ""},
                    {"slide_number": 4, "title": "Next Steps", "type": "summary", "key_points": [], "notes": ""},
                ],
                "data_needed": [],
                "next_steps": []
            }
        try:
            result["pptx_base64"] = self._generate_pptx(result)
            safe_title = result.get("title", "presentation").replace(" ", "_").lower()
            result["pptx_filename"] = f"{safe_title}.pptx"
        except Exception:
            result["pptx_base64"] = ""
            result["pptx_filename"] = ""
        return result

    def _generate_pptx(self, outline: dict) -> str:
        import base64
        import io
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        prs.slide_width = 12192000   # 13.33 inches in EMU
        prs.slide_height = 6858000   # 7.5 inches in EMU

        slides_data = outline.get("slides", [])
        for i, slide_data in enumerate(slides_data):
            layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = slide_data.get("title", f"Slide {i + 1}")
                if title_shape.text_frame.paragraphs:
                    for run in title_shape.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(32 if i == 0 else 24)

            key_points = slide_data.get("key_points", [])
            notes_text = slide_data.get("notes", "")

            if i == 0:
                if len(slide.placeholders) > 1:
                    subtitle = slide.placeholders[1]
                    subtitle.text = outline.get("objective", "") or (key_points[0] if key_points else "")
            else:
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for j, point in enumerate(key_points):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = str(point)
                        p.level = 0
                        for run in p.runs:
                            run.font.size = Pt(16)

            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return base64.b64encode(buf.read()).decode("utf-8")
